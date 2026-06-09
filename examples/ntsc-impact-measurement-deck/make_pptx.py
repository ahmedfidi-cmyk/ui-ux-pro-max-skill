#!/usr/bin/env python3
"""Generate an editable, on-brand NTSC PPTX (RTL Arabic) mirroring the redesign."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- brand palette ----
EVENING = RGBColor(0x00,0x4E,0x43)
GABLE   = RGBColor(0x17,0x33,0x2F)
HIGHLAND= RGBColor(0x71,0x94,0x6A)
SWIRL   = RGBColor(0xD5,0xD0,0xC3)
ANZAC   = RGBColor(0xDE,0xB8,0x3B)
ANZAC_S = RGBColor(0xE9,0xCE,0x74)
PAPER   = RGBColor(0xF6,0xF4,0xEE)
PAPER2  = RGBColor(0xFB,0xFA,0xF6)
INK     = RGBColor(0x2C,0x3A,0x36)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
CREAM   = RGBColor(0xCF,0xE0,0xD9)
AR = "Segoe UI"   # guaranteed Arabic-capable; swap to brand "TheSans" if installed

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def _set_rtl(para, rtl=True, align=PP_ALIGN.RIGHT):
    para.alignment = align
    pPr = para._p.get_or_add_pPr()
    pPr.set('rtl', '1' if rtl else '0')

def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.RIGHT,
         rtl=True, font=AR, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: str or list of (txt,{overrides})"""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs, str): runs=[runs]
    first=True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.line_spacing=line_spacing
        if space_after: p.space_after=Pt(space_after)
        if isinstance(item, tuple): txt, ov = item
        else: txt, ov = item, {}
        _set_rtl(p, ov.get('rtl', rtl), ov.get('align', align))
        r = p.add_run(); r.text = txt
        f = r.font
        f.size = Pt(ov.get('size', size)); f.bold = ov.get('bold', bold)
        f.name = ov.get('font', font)
        f.color.rgb = ov.get('color', color)
        # set complex-script font too
        rPr = r._r.get_or_add_rPr()
        cs = rPr.makeelement(qn('a:cs'), {'typeface': ov.get('font', font)})
        rPr.append(cs)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08, shadow=False, grad=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius>0 else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, x, y, w, h)
    if radius>0:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if grad:
        _gradient(shp, grad[0], grad[1])
    elif fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow: _soft_shadow(shp)
    return shp

def oval(s, x, y, w, h, fill=None, line=None, line_w=1.5, grad=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if grad: _gradient(shp, grad[0], grad[1])
    elif fill is not None: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    else: shp.fill.background()
    if line is not None: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    else: shp.line.fill.background()
    shp.shadow.inherit=False
    return shp

def _gradient(shape, c1, c2, angle=90):
    sp = shape.fill._xPr
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e=sp.find(qn(tag))
        if e is not None: sp.remove(e)
    g = sp.makeelement(qn('a:gradFill'), {})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos,col in ((0,c1),(100000,c2)):
        gs=g.makeelement(qn('a:gs'),{'pos':str(pos)})
        sc=g.makeelement(qn('a:srgbClr'),{'val':'%02X%02X%02X'%(col[0],col[1],col[2])})
        gs.append(sc); lst.append(gs)
    g.append(lst)
    lin=g.makeelement(qn('a:lin'),{'ang':str(angle*60000),'scaled':'1'})
    g.append(lin)
    # insert gradient before line element if present
    ln = sp.find(qn('a:ln'))
    if ln is not None: ln.addprevious(g)
    else: sp.append(g)

def _soft_shadow(shape):
    spPr = shape._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
        {'blurRad':'180000','dist':'60000','dir':'5400000','rotWithShape':'0'})
    c = sh.makeelement(qn('a:srgbClr'), {'val':'07281F'})
    a = c.makeelement(qn('a:alpha'), {'val':'30000'}); c.append(a); sh.append(c); el.append(sh)
    spPr.append(el)

def line_h(s, x, y, w, color=HIGHLAND, weight=1.0):
    ln = s.shapes.add_connector(2, x, y, x+w, y)
    ln.line.color.rgb=color; ln.line.width=Pt(weight); return ln

def pic(s, path, x, y, w=None, h=None):
    return s.shapes.add_picture(path, x, y, width=w, height=h)

def topbar(s, ar, en, num):
    line_h(s, Inches(11.7), Inches(0.62), Inches(0.55), ANZAC, 2)
    text(s, Inches(8.0), Inches(0.40), Inches(4.0), Inches(0.7),
         [(ar,{'size':13,'bold':True,'color':EVENING}), (en,{'size':9,'color':HIGHLAND})],
         align=PP_ALIGN.RIGHT, line_spacing=1.05)
    text(s, Inches(0.6), Inches(0.45), Inches(2.5), Inches(0.4),
         [(num,{'size':11,'bold':True,'color':HIGHLAND,'font':'Consolas','rtl':False,'align':PP_ALIGN.LEFT})],
         rtl=False, align=PP_ALIGN.LEFT)

def footer(s):
    line_h(s, Inches(0.6), Inches(6.95), Inches(12.13), RGBColor(0xCF,0xCB,0xBE), 0.75)
    text(s, Inches(6.8), Inches(7.02), Inches(5.93), Inches(0.35),
         [("المركز الوطني لسلامة النقل · إدارة التواصل",{'size':10.5,'color':HIGHLAND})],
         align=PP_ALIGN.RIGHT)
    text(s, Inches(0.6), Inches(7.05), Inches(5.0), Inches(0.3),
         [("NTSC — IMPACT MEASUREMENT MODEL",{'size':9,'color':HIGHLAND,'rtl':False,'align':PP_ALIGN.LEFT,'font':'Consolas'})],
         rtl=False, align=PP_ALIGN.LEFT)

def title(s, parts, y=Inches(1.05)):
    # parts: list of (txt,color)
    runs=[(t,{'size':30,'bold':True,'color':c}) for t,c in parts]
    # combine on one line: build single paragraph
    tb=s.shapes.add_textbox(Inches(0.6),y,Inches(12.13),Inches(0.8)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; _set_rtl(p,True,PP_ALIGN.RIGHT)
    for t,c in parts:
        r=p.add_run(); r.text=t; f=r.font; f.size=Pt(30);f.bold=True;f.name=AR;f.color.rgb=c
        rPr=r._r.get_or_add_rPr(); rPr.append(rPr.makeelement(qn('a:cs'),{'typeface':AR}))
    return tb

# ============================== S1 COVER ==============================
s=slide(); bg(s, GABLE)
rect(s,0,0,SW,SH,grad=(GABLE,EVENING))
# subtle emblem watermark
p=pic(s,"assets/ntsc-emblem-gold.png", Inches(9.3), Inches(3.6), w=Inches(5.2))
p.element.spPr.append(p.element.spPr.makeelement(qn('a:effectLst'),{}))  # noop
# transparency on watermark
from pptx.oxml.ns import nsmap
blip=p.element.blipFill.find(qn('a:blip'))
am=blip.makeelement(qn('a:alphaModFix'),{'amt':'7000'}); blip.append(am)
pic(s,"assets/ntsc-logo-white.png", Inches(8.1), Inches(0.7), w=Inches(4.4))
text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.7),
     [("نموذج قياس أثر المشاركات",{'size':46,'bold':True,'color':WHITE}),
      ("في الفعاليات والمؤتمرات",{'size':46,'bold':True,'color':WHITE})],
     align=PP_ALIGN.RIGHT, line_spacing=1.08)
# gold accent on middle words via separate box overlay isn't trivial; keep clean
text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.5),
     [("PARTICIPATION IMPACT MEASUREMENT MODEL",{'size':16,'bold':False,'color':ANZAC_S,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],
     rtl=False, align=PP_ALIGN.RIGHT)
rect(s, Inches(11.78), Inches(4.85), Inches(0.95), Inches(0.05), fill=ANZAC, radius=0)
text(s, Inches(6.0), Inches(5.05), Inches(6.73), Inches(0.5),
     [("إدارة التواصل  |  Communications Department",{'size':15,'color':CREAM})],
     align=PP_ALIGN.RIGHT)
text(s, Inches(0.9), Inches(6.7), Inches(7.0), Inches(0.4),
     [("المركز الوطني لسلامة النقل · National Transport Safety Center",{'size':11,'color':RGBColor(0x9F,0xB7,0xB0),'rtl':False,'align':PP_ALIGN.LEFT})],
     rtl=False, align=PP_ALIGN.LEFT)

# ============================== S2 OVERVIEW ==============================
s=slide(); bg(s, PAPER)
topbar(s,"الإطار العام","The Mechanism","01 — 07")
title(s,[("آلية قياس ",EVENING),("أثر المشاركة",ANZAC)])
text(s, Inches(0.6), Inches(1.75), Inches(12.13), Inches(0.7),
     [("منظومة مؤسسية لقياس فاعلية مشاركات المركز، تعتمد على منهجية ومعايير محددة تحوّل الحضور في الفعاليات إلى أثرٍ قابلٍ للقياس وقرارات قائمة على البيانات.",{'size':14.5,'color':INK})],
     align=PP_ALIGN.RIGHT, line_spacing=1.4)
cards=[("01","قياس الأثر الفعلي","تمكين الإدارة من قياس الأثر الحقيقي للمشاركة ومقارنة المُستهدف بالناتج الفعلي.","MEASURE REAL IMPACT"),
       ("02","توثيق الدروس والتحسين","توثيق الدروس المستفادة بشكلٍ منهجي ووضع خطط التحسين المناسبة للمشاركات القادمة.","DOCUMENT & IMPROVE"),
       ("03","دعم التسويق الاستراتيجي","تعزيز الصورة الذهنية للمركز ودعم التسويق الاستراتيجي لحضوره المؤسسي.","STRATEGIC POSITIONING")]
cw=Inches(3.85); gap=Inches(0.29); x0=Inches(0.6); y=Inches(2.75); ch=Inches(3.7)
# RTL order: card 01 on right
for i,(idx,h,b,en) in enumerate(cards):
    x = SW - x0 - cw - i*(cw+gap)
    c=rect(s,x,y,cw,ch,fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.06,shadow=True)
    rect(s,x,y,cw,Inches(0.07),fill=ANZAC,radius=0)
    ic=rect(s,x+cw-Inches(1.0),y+Inches(0.32),Inches(0.7),Inches(0.7),grad=(EVENING,GABLE),radius=0.25)
    text(s,x+Inches(0.3),y+Inches(0.4),Inches(2.4),Inches(0.4),[(idx,{'size':12,'bold':True,'color':ANZAC,'rtl':False,'align':PP_ALIGN.LEFT,'font':'Consolas'})],rtl=False,align=PP_ALIGN.LEFT)
    text(s,x+Inches(0.28),y+Inches(1.25),cw-Inches(0.56),Inches(0.6),[(h,{'size':17,'bold':True,'color':EVENING})],align=PP_ALIGN.RIGHT)
    text(s,x+Inches(0.28),y+Inches(1.95),cw-Inches(0.56),Inches(1.2),[(b,{'size':12.5,'color':RGBColor(0x52,0x61,0x5C)})],align=PP_ALIGN.RIGHT,line_spacing=1.35)
    text(s,x+Inches(0.28),y+ch-Inches(0.55),cw-Inches(0.56),Inches(0.35),[(en,{'size':9.5,'color':HIGHLAND,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.RIGHT)
footer(s)

# ============================== S3 METHODOLOGY ==============================
s=slide(); bg(s, PAPER)
topbar(s,"المنهجية العامة","The Methodology","02 — 07")
title(s,[("منهجية القياس في ",EVENING),("خمس مراحل",ANZAC)])
text(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(0.5),
     [("مسارٌ متكامل ينقل المشاركة من جمع البيانات الأولية إلى أثرٍ موثّق وخطط تحسين قابلة للتنفيذ.",{'size':14,'color':INK})],
     align=PP_ALIGN.RIGHT)
steps=[("1","جمع المعلومات","DATA COLLECTION"),("2","تحديد المخرجات","DEFINE OUTPUTS"),
       ("3","قياس المؤشرات","MEASURE KPIS"),("4","الأثر والدروس","IMPACT & INSIGHTS"),
       ("5","خطط التحسين","OPTIMIZATION")]
n=5; node=Inches(1.15); area_x=Inches(0.7); area_w=SW-Inches(1.4)
ystep=Inches(2.95)
slot=(area_w-node)/(n-1)
cy=ystep+node/2
line_h(s, area_x+node/2, cy, area_w-node, HIGHLAND, 1.25)
for i,(num,h,en) in enumerate(steps):
    # RTL: step1 on right
    cx = SW - area_x - node - i*slot
    on = (i==0)
    if on: oval(s,cx,ystep,node,node,grad=(EVENING,GABLE))
    else: oval(s,cx,ystep,node,node,fill=PAPER,line=RGBColor(0xCF,0xDA,0xD2),line_w=2)
    text(s,cx,ystep+Inches(0.33),node,Inches(0.5),[(num,{'size':22,'bold':True,'color':(ANZAC if on else EVENING),'rtl':False,'align':PP_ALIGN.CENTER})],rtl=False,align=PP_ALIGN.CENTER)
    badge=oval(s,cx+node-Inches(0.36),ystep-Inches(0.05),Inches(0.42),Inches(0.42),fill=ANZAC)
    text(s,cx+node-Inches(0.36),ystep+Inches(0.01),Inches(0.42),Inches(0.32),[(num,{'size':12,'bold':True,'color':GABLE,'rtl':False,'align':PP_ALIGN.CENTER})],rtl=False,align=PP_ALIGN.CENTER)
    text(s,cx-Inches(0.5),ystep+node+Inches(0.15),node+Inches(1.0),Inches(0.4),[(h,{'size':14,'bold':True,'color':EVENING})],align=PP_ALIGN.CENTER)
    text(s,cx-Inches(0.5),ystep+node+Inches(0.55),node+Inches(1.0),Inches(0.3),[(en,{'size':8.5,'color':HIGHLAND,'bold':True,'rtl':False,'align':PP_ALIGN.CENTER,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.CENTER)
# phase chips
chips=["قبل المشاركة · التخطيط","أثناء وبعد المشاركة · التنفيذ والقياس","ما بعد المشاركة · التحسين"]
cwd=Inches(3.4); cy2=Inches(5.55)
for i,t in enumerate(chips):
    x=SW-Inches(0.9)-cwd-i*(cwd+Inches(0.25))
    rect(s,x,cy2,cwd,Inches(0.55),fill=RGBColor(0xEA,0xEE,0xE7),line=RGBColor(0xCF,0xDA,0xD2),line_w=0.75,radius=0.5)
    text(s,x,cy2+Inches(0.1),cwd,Inches(0.4),[(t,{'size':11.5,'color':EVENING})],align=PP_ALIGN.CENTER)
footer(s)

# ============================== S4 DATA COLLECTION ==============================
s=slide(); bg(s, PAPER)
topbar(s,"المرحلة الأولى","Phase 01 · Data Collection","03 — 07")
title(s,[("جمع ",EVENING),("المعلومات",ANZAC)])
pw=Inches(5.95); ph=Inches(3.55); py=Inches(2.0)
# right panel = before (light)
xr=SW-Inches(0.6)-pw
rect(s,xr,py,pw,ph,fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.05,shadow=True)
text(s,xr+Inches(0.4),py+Inches(0.3),pw-Inches(0.8),Inches(0.6),[("قبل الحدث",{'size':18,'bold':True,'color':EVENING}),("BEFORE THE EVENT",{'size':9.5,'color':HIGHLAND,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],align=PP_ALIGN.RIGHT)
before=["عدد الحضور المتوقع للفعالية","نوع الجمهور المستهدَف وطبيعته","مشاركة الجهات ذات العلاقة","نوعية التغطية الإعلامية المتوقعة"]
text(s,xr+Inches(0.4),py+Inches(1.25),pw-Inches(0.8),Inches(2.1),
     [("◆  "+t,{'size':14,'color':RGBColor(0x3F,0x4F,0x4A)}) for t in before],align=PP_ALIGN.RIGHT,line_spacing=1.3,space_after=10)
# left panel = after (dark)
xl=Inches(0.6)
rect(s,xl,py,pw,ph,grad=(EVENING,GABLE),radius=0.05,shadow=True)
text(s,xl+Inches(0.4),py+Inches(0.3),pw-Inches(0.8),Inches(0.6),[("خلال وبعد الحدث",{'size':18,'bold':True,'color':WHITE}),("DURING & AFTER",{'size':9.5,'color':ANZAC_S,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],align=PP_ALIGN.RIGHT)
after=["عدد زوّار الجناح الفعلي","تحليل عام لفئات زوّار الجناح","العدد الفعلي للجهات المشاركة ذات العلاقة","عدد التغطيات الإعلامية الخارجية للمركز"]
text(s,xl+Inches(0.4),py+Inches(1.25),pw-Inches(0.8),Inches(2.1),
     [("◆  "+t,{'size':14,'color':RGBColor(0xE6,0xEE,0xEA)}) for t in after],align=PP_ALIGN.RIGHT,line_spacing=1.3,space_after=10)
# note
rect(s,Inches(0.6),Inches(5.85),Inches(12.13),Inches(0.7),fill=RGBColor(0xEC,0xF0,0xE9),line=None,radius=0.1)
rect(s,Inches(12.66),Inches(5.85),Inches(0.07),Inches(0.7),fill=ANZAC,radius=0)
text(s,Inches(0.9),Inches(6.02),Inches(11.5),Inches(0.4),[("يتم جمع المعلومات يدويًا عن طريق فريق ممثّلي المركز خلال المشاركة.",{'size':12.5,'color':EVENING})],align=PP_ALIGN.RIGHT)
footer(s)

# ============================== S5 KPIs ==============================
s=slide(); bg(s, PAPER)
topbar(s,"المرحلتان الثانية والثالثة","Phase 02–03 · Outputs & KPIs","04 — 07")
title(s,[("المخرجات وقياس ",EVENING),("المؤشرات",ANZAC)])
kpis=[("عدد زوّار الجناح","٥٪ من الزوّار (٥١ جناحًا أو أكثر) · ١٠٪ (٥٠ أو أقل)","5–10%","من الزوّار","RELATIVE · نسبي",ANZAC),
      ("مدّة بقاء الزائر","متوسط الزمن داخل جناح المركز","2–5","دقائق","QUANT · كمّي",EVENING),
      ("الشراكات المستهدَفة","عدد الشراكات النوعية المستهدَفة","2","شراكات","QUANT · كمّي",EVENING),
      ("التغطيات الإعلامية الخارجية","تغطيات من جهات خارج المركز","1","تغطية","QUANT · كمّي",EVENING)]
lx=Inches(4.7); lw=Inches(8.05); ry=Inches(1.95); rh=Inches(1.07); rgap=Inches(0.16)
for i,(h,d,v,u,tg,tc) in enumerate(kpis):
    y=ry+i*(rh+rgap)
    rect(s,lx,y,lw,rh,fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.1,shadow=True)
    # value left
    text(s,lx+Inches(0.25),y+Inches(0.2),Inches(1.4),Inches(0.55),[(v,{'size':23,'bold':True,'color':EVENING,'rtl':False,'align':PP_ALIGN.CENTER})],rtl=False,align=PP_ALIGN.CENTER)
    text(s,lx+Inches(0.25),y+Inches(0.72),Inches(1.4),Inches(0.3),[(u,{'size':10,'color':HIGHLAND,'align':PP_ALIGN.CENTER})],align=PP_ALIGN.CENTER)
    # title + desc right
    text(s,lx+Inches(2.0),y+Inches(0.16),lw-Inches(3.2),Inches(0.4),[(h,{'size':15.5,'bold':True,'color':EVENING})],align=PP_ALIGN.RIGHT)
    text(s,lx+Inches(2.0),y+Inches(0.56),lw-Inches(3.2),Inches(0.45),[(d,{'size':10.5,'color':RGBColor(0x6A,0x78,0x73)})],align=PP_ALIGN.RIGHT,line_spacing=1.15)
    # icon box right end
    rect(s,lx+lw-Inches(0.95),y+Inches(0.24),Inches(0.6),Inches(0.6),grad=(EVENING,GABLE),radius=0.22)
# side gauge panel
sx=Inches(0.6); sw=Inches(3.95)
rect(s,sx,ry,sw,Inches(4.6),grad=(EVENING,GABLE),radius=0.06,shadow=True)
text(s,sx+Inches(0.35),ry+Inches(0.3),sw-Inches(0.7),Inches(0.6),[("قياس الأداء",{'size':15,'bold':True,'color':WHITE}),("TARGET VS. ACTUAL",{'size':9.5,'color':ANZAC_S,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],align=PP_ALIGN.RIGHT)
# gauge ring (gold oval + inner dark)
gx=sx+sw/2-Inches(0.9)
oval(s,gx,ry+Inches(1.15),Inches(1.8),Inches(1.8),fill=ANZAC)
oval(s,gx+Inches(0.22),ry+Inches(1.37),Inches(1.36),Inches(1.36),fill=GABLE)
text(s,gx,ry+Inches(1.65),Inches(1.8),Inches(0.6),[("100%",{'size':26,'bold':True,'color':WHITE,'rtl':False,'align':PP_ALIGN.CENTER})],rtl=False,align=PP_ALIGN.CENTER)
text(s,gx-Inches(0.4),ry+Inches(2.25),Inches(2.6),Inches(0.3),[("تحقيق المستهدف",{'size':10,'color':ANZAC_S,'align':PP_ALIGN.CENTER})],align=PP_ALIGN.CENTER)
text(s,sx+Inches(0.35),ry+Inches(3.4),sw-Inches(0.7),Inches(1.0),[("تُحوَّل كل مخرجات المشاركة إلى مؤشرات نسبية أو كمّية لمقارنة الأداء الفعلي بالمستهدف.",{'size':11,'color':CREAM})],align=PP_ALIGN.RIGHT,line_spacing=1.3)
footer(s)

# ============================== S6 IMPACT & OPTIMIZATION ==============================
s=slide(); bg(s, PAPER)
topbar(s,"المرحلتان الرابعة والخامسة","Phase 04–05 · Impact & Optimization","05 — 07")
title(s,[("الأثر والدروس ",EVENING),("وخطط التحسين",ANZAC)])
blocks=[("4","الأثر والدروس المستفادة","IMPACTS & INSIGHTS",
         ["تحليل الأداء العام ومقارنة المستهدَف بالناتج الفعلي.","تحديد نقاط القوة ووضع خطة لتطويرها.","تحديد التحديات ووضع خطة لمعالجتها."]),
        ("5","خطط التحسين","OPTIMIZATION PLANS",
         ["تقديم التقرير النهائي لمدير عام الإدارة.","وضع خطط لمعالجة التحديات وتحسين المشاركات المستقبلية.","تحديد الميزانيات المطلوبة (إن وُجدت)."])]
bw=Inches(5.95); bh=Inches(3.9); by=Inches(2.0)
for i,(num,h,en,items) in enumerate(blocks):
    x = (SW-Inches(0.6)-bw) if i==0 else Inches(0.6)
    rect(s,x,by,bw,bh,fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.05,shadow=True)
    rect(s,x+bw-Inches(1.0),by+Inches(0.35),Inches(0.65),Inches(0.65),grad=(ANZAC,RGBColor(0xC7,0x9A,0x26)),radius=0.22)
    text(s,x+bw-Inches(1.0),by+Inches(0.42),Inches(0.65),Inches(0.5),[(num,{'size':20,'bold':True,'color':GABLE,'rtl':False,'align':PP_ALIGN.CENTER})],rtl=False,align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.4),by+Inches(0.4),bw-Inches(1.6),Inches(0.4),[(h,{'size':17,'bold':True,'color':EVENING})],align=PP_ALIGN.RIGHT)
    text(s,x+Inches(0.4),by+Inches(0.85),bw-Inches(1.6),Inches(0.3),[(en,{'size':9.5,'color':HIGHLAND,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.RIGHT)
    line_h(s,x+Inches(0.4),by+Inches(1.35),bw-Inches(0.8),RGBColor(0xDD,0xE4,0xDF),1)
    text(s,x+Inches(0.4),by+Inches(1.6),bw-Inches(0.8),Inches(2.1),
         [("◆  "+t,{'size':13.5,'color':RGBColor(0x3F,0x4F,0x4A)}) for t in items],align=PP_ALIGN.RIGHT,line_spacing=1.3,space_after=12)
footer(s)

# ============================== S7 DOCUMENTATION ==============================
s=slide(); bg(s, PAPER)
topbar(s,"الأرشفة المؤسسية","Knowledge Management","06 — 07")
title(s,[("إدارة المعلومات ",EVENING),("والتوثيق",ANZAC)])
# right text + cadence
text(s,Inches(6.9),Inches(2.0),Inches(5.83),Inches(1.4),
     [("تُحفظ التقارير والدراسات التسويقية المُعدّة قبل كل مشاركة، وتُجمَّع في أرشيفٍ متكامل يُمكّن المركز من إجراء تحليلات تسويقية شاملة عبر أطرٍ زمنية أطول.",{'size':14,'color':INK})],
     align=PP_ALIGN.RIGHT,line_spacing=1.45)
cad=[("سنوي","ANNUAL"),("ربع سنوي","QUARTERLY"),("موسمي","SEASONAL")]
cwid=Inches(1.83); cyy=Inches(3.7)
for i,(a,e) in enumerate(cad):
    x=SW-Inches(0.6)-cwid-i*(cwid+Inches(0.17))
    rect(s,x,cyy,cwid,Inches(1.5),fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.08,shadow=True)
    oval(s,x+cwid/2-Inches(0.3),cyy+Inches(0.22),Inches(0.6),Inches(0.6),grad=(EVENING,GABLE))
    text(s,x,cyy+Inches(0.95),cwid,Inches(0.35),[(a,{'size':14,'bold':True,'color':EVENING,'align':PP_ALIGN.CENTER})],align=PP_ALIGN.CENTER)
    text(s,x,cyy+Inches(1.2),cwid,Inches(0.25),[(e,{'size':8.5,'color':HIGHLAND,'bold':True,'rtl':False,'align':PP_ALIGN.CENTER,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.CENTER)
# left archive flow panel (dark)
ax=Inches(0.6); aw=Inches(5.95)
rect(s,ax,Inches(2.0),aw,Inches(4.5),grad=(EVENING,GABLE),radius=0.06,shadow=True)
flow=[("تقارير ما بعد المشاركة","POST-PARTICIPATION REPORTS"),
      ("الدراسات التسويقية المسبقة","PRE-EVENT MARKETING STUDIES"),
      ("أرشيف بيانات موحّد","UNIFIED DATA ARCHIVE"),
      ("تحليلات تسويقية شاملة","COMPREHENSIVE ANALYTICS")]
for i,(a,e) in enumerate(flow):
    yy=Inches(2.3)+i*Inches(1.02)
    rect(s,ax+aw-Inches(0.95),yy+Inches(0.12),Inches(0.6),Inches(0.6),fill=RGBColor(0x0A,0x3A,0x33),radius=0.22)
    text(s,ax+Inches(0.4),yy+Inches(0.12),aw-Inches(1.5),Inches(0.4),[(a,{'size':14,'bold':True,'color':WHITE})],align=PP_ALIGN.RIGHT)
    text(s,ax+Inches(0.4),yy+Inches(0.5),aw-Inches(1.5),Inches(0.3),[(e,{'size':8.5,'color':ANZAC_S,'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.RIGHT)
    if i<3: line_h(s,ax+Inches(0.4),yy+Inches(0.95),aw-Inches(0.8),RGBColor(0x2C,0x52,0x4B),0.75)
footer(s)

# ============================== S8 REPORT TEMPLATE ==============================
s=slide(); bg(s, PAPER)
topbar(s,"المُخرَج النهائي","The Deliverable","07 — 07")
title(s,[("تقرير ",EVENING),("ما بعد المشاركة",ANZAC)])
text(s,Inches(0.6),Inches(1.7),Inches(12.13),Inches(0.5),
     [("نموذجٌ موحّد يجمع نتائج المشاركة في وثيقةٍ واحدة تُرفع لمدير عام الإدارة وتُحفظ في الأرشيف المؤسسي.",{'size':13.5,'color':INK})],align=PP_ALIGN.RIGHT)
secs=[("01","بيانات المشاركة","PARTICIPATION DATA",["اسم الفعالية والجهة المنظِّمة","مكان وتاريخ المشاركة","نوع المشاركة","ممثلو المركز المشاركون"]),
      ("02","ملخص المشاركة","SUMMARY",["نبذة مختصرة عن الفعالية","أهم الأنشطة التي نفّذها المركز","أبرز محطات الحضور"]),
      ("03","النتائج والمؤشرات","RESULTS & KPIS",["عدد زوّار الجناح","متوسط مدّة بقاء الزائر","الشراكات المستهدَفة","التغطيات الإعلامية الخارجية"]),
      ("04","الدروس المستفادة","LESSONS LEARNED",["نقاط القوة","التحديات"]),
      ("05","التوصيات وخطة التحسين","RECOMMENDATIONS",["التوصيات التنفيذية","خطة التحسين للمشاركات القادمة","الميزانيات المطلوبة (إن وُجدت)"]),
      ("✓","جاهز للأرشفة","ARCHIVE-READY",["وثيقة معتمدة تُغذّي التحليلات السنوية","وتدعم اتخاذ القرار."])]
gw=Inches(3.95); gh=Inches(1.95); gx0=Inches(0.6); gy0=Inches(2.35); ggap=Inches(0.14)
for i,(num,h,en,items) in enumerate(secs):
    col=i%3; row=i//3
    x=SW-gx0-gw-col*(gw+ggap)
    y=gy0+row*(gh+ggap)
    dark=(num=="✓")
    if dark: rect(s,x,y,gw,gh,grad=(EVENING,GABLE),radius=0.06,shadow=True)
    else: rect(s,x,y,gw,gh,fill=PAPER2,line=RGBColor(0xCF,0xDA,0xD2),line_w=1,radius=0.06,shadow=True)
    text(s,x+Inches(0.28),y+Inches(0.2),gw-Inches(0.56),Inches(0.35),[(num,{'size':12,'bold':True,'color':(ANZAC_S if dark else ANZAC),'rtl':False,'align':PP_ALIGN.LEFT,'font':'Consolas'})],rtl=False,align=PP_ALIGN.LEFT)
    text(s,x+Inches(0.28),y+Inches(0.18),gw-Inches(0.56),Inches(0.35),[(h,{'size':15,'bold':True,'color':(WHITE if dark else EVENING)})],align=PP_ALIGN.RIGHT)
    text(s,x+Inches(0.28),y+Inches(0.55),gw-Inches(0.56),Inches(0.25),[(en,{'size':8.5,'color':(ANZAC_S if dark else HIGHLAND),'bold':True,'rtl':False,'align':PP_ALIGN.RIGHT,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.RIGHT)
    col_txt = CREAM if dark else RGBColor(0x52,0x61,0x5C)
    text(s,x+Inches(0.28),y+Inches(0.85),gw-Inches(0.56),Inches(1.0),[("•  "+t,{'size':11,'color':col_txt}) for t in items],align=PP_ALIGN.RIGHT,line_spacing=1.15,space_after=3)
footer(s)

# ============================== S9 CLOSING ==============================
s=slide(); bg(s, GABLE)
rect(s,0,0,SW,SH,grad=(GABLE,EVENING))
wm=pic(s,"assets/ntsc-emblem-gold.png", Inches(9.5), Inches(3.8), w=Inches(5.0))
blip=wm.element.blipFill.find(qn('a:blip')); blip.append(blip.makeelement(qn('a:alphaModFix'),{'amt':'6000'}))
pic(s,"assets/ntsc-logo-white.png", Inches(4.97), Inches(1.7), w=Inches(3.4))
text(s,Inches(0),Inches(2.95),SW,Inches(1.3),[("شـكـراً لـكـم",{'size':60,'bold':True,'color':WHITE,'align':PP_ALIGN.CENTER})],align=PP_ALIGN.CENTER)
rect(s,SW/2-Inches(0.5),Inches(4.35),Inches(1.0),Inches(0.05),fill=ANZAC,radius=0)
text(s,Inches(0),Inches(4.6),SW,Inches(0.5),[("نقيس الأثر · نوثّق الدروس · نرتقي بالحضور",{'size':17,'color':CREAM,'align':PP_ALIGN.CENTER})],align=PP_ALIGN.CENTER)
text(s,Inches(0),Inches(5.25),SW,Inches(0.4),[("Communications Department — National Transport Safety Center",{'size':12,'color':HIGHLAND,'rtl':False,'align':PP_ALIGN.CENTER,'font':'Segoe UI'})],rtl=False,align=PP_ALIGN.CENTER)

prs.save("NTSC-Impact-Measurement-Model.pptx")
print("Saved NTSC-Impact-Measurement-Model.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
